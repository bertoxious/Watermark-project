
import os
from wand.image import Image
from pptx import Presentation
from pptx.util import Inches

## Enter url of the folder where images are
## example :: e:\folder\images 
src = input("Enter src folder :: ")
list_of_images = os.listdir(src)

# Enter url of watermark along with watermark name
## example :: e:\folder\watermark.png
watermark = input("Enter watermark url :: ")

# this will set up a frame to work on for our presentation
prs = Presentation()

# here i chose a blank layout
slide_layout = prs.slide_layouts[1]


'''selecting each image from the directory and performing operation
 on each one of it'''
for each_image in list_of_images:
	
	# adding slide or page to ppt
	with Image(filename = f'{src}\\{each_image}') as img:
		img.transform(resize='300x320>')
		print(img.size)

		with Image(filename=watermark) as img2:
			img2.transform(resize='80x80>')
			with img.clone() as demo:
				demo.watermark(img2, 0, 0,0)
				demo.save(filename="watermark.jpg")

				slide = prs.slides.add_slide(slide_layout)

				# adding title
				title = slide.shapes.title
				title.text = f"Sample Title {list_of_images.index(each_image)+1}"
				# title.text_frame.paragraphs[0].font.size = Pt(48)

				# adding subtitle
				subtitle = slide.shapes.placeholders[1]
				subtitle.text = f"Sample subtitle {list_of_images.index(each_image)+1}"

				# adding picture
				pic = slide.shapes.add_picture("watermark.jpg", Inches(0.5), Inches(2.5))#

	os.remove("watermark.jpg")

prs.save("Watermark.pptx")


