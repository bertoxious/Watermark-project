from pptx import Presentation 
from pptx.util import Inches

prs = Presentation()

slide1_register = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide1_register)

title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]


title1.text = "Here comes the header"
subtitle1.text = "Here comes the subtitle"

slide2_register = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(slide2_register)

title2 = slide2.shapes.title 
title2.text = "Here comes the bullet points"

bullet_points = slide2.shapes

bplvl1 = bullet_points.placeholders[1]
bplvl1.text = "Ye aaya pehla"

image = "nike.png"

picture = slide2.shapes.add_picture(image, Inches(3), Inches(3),Inches(5))

slide3_register = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(slide3_register)

title3 = slide3.shapes.title
title3.text = "Adding Image" 

image = "nike.png"

left = Inches(3)
top = Inches(2)

add_picture = slide3.shapes.add_picture(image, left, top)

prs.save("New.pptx")